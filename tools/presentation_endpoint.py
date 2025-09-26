import io
import os
from fastapi import APIRouter, HTTPException
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

router = APIRouter()

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

# Base folder for all templates
TEMPLATES_DIR = os.path.join(os.path.dirname(__file__), "templates")


def load_template(template_name: str | None) -> Presentation:
    """Load a specific template if provided, otherwise return a new blank presentation."""
    if template_name:
        template_path = os.path.join(TEMPLATES_DIR, template_name)
        if os.path.exists(template_path):
            return Presentation(template_path)
        else:
            raise FileNotFoundError(f"Template '{template_name}' not found in {TEMPLATES_DIR}")
    return Presentation()


def generate_ppt(data: dict) -> bytes:
 
    template_name = data.get("template")
    prs = load_template(template_name)

    for i, slide_data in enumerate(data["slides"]):
        if i < len(prs.slides):
            slide = prs.slides[i]
        else:
            layout_index = slide_data.get("layout_index", 1)
            slide_layout = (
                prs.slide_layouts[layout_index]
                if layout_index < len(prs.slide_layouts)
                else prs.slide_layouts[1]
            )
            slide = prs.slides.add_slide(slide_layout)

        # --- Title ---
        if "title" in slide_data and slide.shapes.title:
            t = slide_data["title"]
            title_shape = slide.shapes.title
            title_shape.text = t.get("text", "")

            if hasattr(title_shape, "text_frame"):
                para = title_shape.text_frame.paragraphs[0]
                run = para.runs[0] if para.runs else para.add_run()
                run.font.size = Pt(t.get("font_size", 32))
                run.font.bold = t.get("bold", True)
                run.font.italic = t.get("italic", False)
                run.font.color.rgb = RGBColor(*t.get("color", (0, 0, 0)))
                para.alignment = ALIGN_MAP.get(t.get("align", "center"), PP_ALIGN.CENTER)

        # --- Content ---
        if slide.placeholders and len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            tf = content_placeholder.text_frame
            tf.clear()

            for c in slide_data.get("content", []):
                p = tf.add_paragraph()
                p.text = c["text"]
                p.font.size = Pt(c.get("font_size", 20))
                p.font.bold = c.get("bold", False)
                p.font.italic = c.get("italic", False)
                p.font.color.rgb = RGBColor(*c.get("color", (0, 0, 0)))
                p.alignment = ALIGN_MAP.get(c.get("align", "left"), PP_ALIGN.LEFT)

    # Save PPTX to memory
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


@router.post("/generate-pptx")
async def create_pptx(request: dict):
   
    try:
        ppt_bytes = generate_ppt(request)
        buffer = io.BytesIO(ppt_bytes)

        return StreamingResponse(
            buffer,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f"attachment; filename={request.get('filename', 'generated.pptx')}"
            },
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
