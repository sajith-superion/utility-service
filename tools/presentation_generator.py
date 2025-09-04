import io
import os
from fastapi import APIRouter, HTTPException
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

router = APIRouter()

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY
}

def generate_ppt(data: dict) -> bytes:
   
    # Load template or start new presentation
    if data.get("template", False) and os.path.exists("template.pptx"):
        prs = Presentation("template.pptx")
    else:
        prs = Presentation()

    for slide_data in data["slides"]:
        
        layout_index = slide_data.get("layout_index", 1)
        slide_layout = (
            prs.slide_layouts[layout_index]
            if layout_index < len(prs.slide_layouts)
            else prs.slide_layouts[1]
        )
        slide = prs.slides.add_slide(slide_layout)

        # Title placeholder
        if "title" in slide_data:
            t = slide_data["title"]
            title_shape = slide.shapes.title
            title_shape.text = t.get("text", "")
            if hasattr(title_shape, "text_frame"):
                run = title_shape.text_frame.paragraphs[0].runs[0]
                run.font.size = Pt(t.get("font_size", 32))
                run.font.bold = t.get("bold", True)
                run.font.italic = t.get("italic", False)
                run.font.color.rgb = RGBColor(*t.get("color", (0, 0, 0)))
                title_shape.text_frame.paragraphs[0].alignment = ALIGN_MAP.get(
                    t.get("align", "center"), PP_ALIGN.CENTER
                )

        # Content placeholder
        if slide.placeholders and len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            tf = content_placeholder.text_frame
            tf.clear()  # start fresh

            for c in slide_data.get("content", []):
                p = tf.add_paragraph()
                p.text = c["text"]
                p.font.size = Pt(c.get("font_size", 20))
                p.font.bold = c.get("bold", False)
                p.font.italic = c.get("italic", False)
                p.font.color.rgb = RGBColor(*c.get("color", (0, 0, 0)))
                p.alignment = ALIGN_MAP.get(c.get("align", "left"), PP_ALIGN.LEFT)

    # Save PPTX into memory
    buffer = io.BytesIO()
    prs.save(buffer)
    ppt_bytes = buffer.getvalue()
    buffer.close()
    return ppt_bytes


@router.post("/generate-pptx")
async def create_pptx(request: dict):
   
    try:
        ppt_bytes = generate_ppt(request)
        buffer = io.BytesIO(ppt_bytes)

        return StreamingResponse(
            buffer,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f"attachment; filename={request.get('filename', 'generated.pptx')}"
            },
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
